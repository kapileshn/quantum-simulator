"""
Generate a PowerPoint presentation for the Quantum Protocol & Algorithm Simulator.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Colors ──
BLUE = RGBColor(0x1a, 0x56, 0xdb)
DARK_BLUE = RGBColor(0x0f, 0x2b, 0x6e)
DARK = RGBColor(0x1f, 0x29, 0x37)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x6b, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0x9c, 0xa3, 0xaf)
ACCENT = RGBColor(0x7c, 0x3a, 0xed)
GREEN = RGBColor(0x05, 0x96, 0x69)
ORANGE = RGBColor(0xd9, 0x77, 0x06)
RED = RGBColor(0xdc, 0x26, 0x26)
BG_DARK = RGBColor(0x0f, 0x17, 0x2a)
BG_SLIDE = RGBColor(0xf8, 0xfa, 0xff)
CARD_BG = RGBColor(0xe8, 0xee, 0xfb)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=DARK,
             bold=False, italic=False, align=PP_ALIGN.LEFT, font="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    return tf


def add_bullet_frame(slide, left, top, width, height, items, font_size=16,
                     color=DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        # bullet char
        p.text = "\u2022  " + item
    return tf


def add_card(slide, left, top, width, height, title_text, body_items,
             accent_color=BLUE, font_size=14):
    # Card background
    add_shape(slide, left, top, width, height, CARD_BG)
    # Accent bar
    add_shape(slide, left, top, Inches(0.06), height, accent_color)
    # Title
    add_text(slide, left + Inches(0.25), top + Inches(0.1), width - Inches(0.4),
             Inches(0.4), title_text, font_size=17, color=accent_color, bold=True)
    # Body
    add_bullet_frame(slide, left + Inches(0.25), top + Inches(0.5),
                     width - Inches(0.4), height - Inches(0.6),
                     body_items, font_size=font_size, color=DARK)


def section_divider(title_text, subtitle_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, BG_DARK)
    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
             title_text, font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if subtitle_text:
        add_text(slide, Inches(1), Inches(4), Inches(11), Inches(1),
                 subtitle_text, font_size=22, color=LIGHT_GRAY, italic=True,
                 align=PP_ALIGN.CENTER)
    return slide


def content_slide(title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_SLIDE)
    # Top bar
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE)
    # Title
    add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             title_text, font_size=32, color=DARK_BLUE, bold=True)
    # Divider line
    add_shape(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.03), BLUE)
    return slide


# ════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_text(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.5),
         "Quantum Protocol &\nAlgorithm Simulator",
         font_size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.6), Inches(11.3), Inches(0.8),
         "An interactive, physics-accurate platform for quantum computing education",
         font_size=24, color=LIGHT_GRAY, italic=True, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.5), Inches(4.6), Inches(2.3), Inches(0.04), BLUE)
add_text(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.5),
         "Version 1.0  \u2022  2026",
         font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ════════════════════════════════════════════════════════════
slide = content_slide("Agenda")
items = [
    "1.  What is a Quantum Simulator?",
    "2.  System Architecture",
    "3.  The Quantum State \u2014 Qubits, Superposition & Hilbert Space",
    "4.  Complete Gate Library (17 gates, all physics-verified)",
    "5.  Measurement & Collapse",
    "6.  Bloch Sphere Visualization",
    "7.  Entanglement & Bell States",
    "8.  Algorithm: Deutsch-Jozsa",
    "9.  Algorithm: Grover\u2019s Search",
    "10. Protocol: Quantum Teleportation",
    "11. Protocol: BB84 Quantum Key Distribution",
    "12. Quantum Random Number Generator",
    "13. Live Demo Walkthrough",
]
add_bullet_frame(slide, Inches(1.5), Inches(1.3), Inches(10), Inches(5.5),
                 items, font_size=19, color=DARK, spacing=Pt(8))

# ════════════════════════════════════════════════════════════
# SLIDE 3: WHAT IS IT
# ════════════════════════════════════════════════════════════
slide = content_slide("What is the Quantum Simulator?")
add_bullet_frame(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(3),
    [
        "Full statevector simulation of 1\u20135 qubit quantum systems",
        "17 physics-verified unitary gates (all satisfy U\u2020U = I)",
        "Step-by-step circuit execution with time-travel scrubbing",
        "Real-time 3D Bloch sphere visualization per qubit",
        "5 pre-built algorithms & cryptographic protocols",
        "WebSocket-powered live state updates",
    ], font_size=17, color=DARK)

add_card(slide, Inches(7.5), Inches(1.3), Inches(5), Inches(2.2),
         "Tech Stack", [
             "Backend: Python \u2022 FastAPI \u2022 NumPy \u2022 Socket.IO",
             "Frontend: React \u2022 TypeScript \u2022 Three.js \u2022 Zustand",
         ], accent_color=ACCENT, font_size=15)

add_card(slide, Inches(7.5), Inches(3.8), Inches(5), Inches(2.2),
         "Key Principle", [
             "Every gate matrix is explicit \u2014 no black-box libraries.",
             "Physics first: all math is transparent and verifiable.",
         ], accent_color=GREEN, font_size=15)

# ════════════════════════════════════════════════════════════
# SLIDE 4: ARCHITECTURE
# ════════════════════════════════════════════════════════════
slide = content_slide("System Architecture")

# Frontend box
add_card(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(2.5),
         "Frontend (React + TypeScript)", [
             "AlgorithmSelector \u2014 configure & run algorithms",
             "CircuitEditor \u2014 drag-and-drop gate placement",
             "BlochSphere \u2014 3D WebGL per-qubit visualization",
             "StatevectorPanel \u2014 amplitude/probability tables",
             "Zustand store \u2014 centralized state management",
         ], accent_color=BLUE, font_size=14)

# Backend box
add_card(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(2.5),
         "Backend (FastAPI + NumPy)", [
             "QuantumState \u2014 statevector engine (2\u207F amplitudes)",
             "Gates module \u2014 17 unitary matrices",
             "Measurement \u2014 projective collapse & renormalization",
             "Algorithms \u2014 D-J, Grover, Teleport, BB84, QRNG",
             "Socket.IO \u2014 real-time step-by-step push",
         ], accent_color=ACCENT, font_size=14)

# Connection arrow
add_text(slide, Inches(4.5), Inches(4.2), Inches(4.5), Inches(0.5),
         "\u2b04  WebSocket / REST API  \u2b04",
         font_size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# API table
add_card(slide, Inches(2), Inches(4.9), Inches(9), Inches(1.8),
         "API Endpoints", [
             "GET  /api/gates \u2014 list all gates     |     GET  /api/algorithms \u2014 list algorithms",
             "POST /api/simulate \u2014 run custom circuit     |     POST /api/algorithms/{name}/run",
             "Socket.IO: start_session \u2192 session_started \u2192 step \u2192 step_update",
         ], accent_color=ORANGE, font_size=14)

# ════════════════════════════════════════════════════════════
# SLIDE 5: QUANTUM STATE
# ════════════════════════════════════════════════════════════
section_divider("The Quantum State",
                "Qubits, Superposition & Hilbert Space")

slide = content_slide("Qubits & Superposition")
add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
         "A qubit is a two-level quantum system. Unlike a classical bit (0 or 1), a qubit exists in superposition:",
         font_size=18, color=DARK)

# Math box
add_shape(slide, Inches(2.5), Inches(2.0), Inches(8), Inches(0.7), CARD_BG)
add_text(slide, Inches(2.7), Inches(2.05), Inches(7.6), Inches(0.6),
         "|\u03c8\u27e9 = \u03b1|0\u27e9 + \u03b2|1\u27e9      where  \u03b1, \u03b2 \u2208 \u2102   and   |\u03b1|\u00b2 + |\u03b2|\u00b2 = 1",
         font_size=22, color=DARK, font="Cambria Math", align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(2.9), Inches(11), Inches(0.5),
         "For n qubits, the state lives in a 2\u207f-dimensional complex Hilbert space:",
         font_size=18, color=DARK)

add_shape(slide, Inches(2.5), Inches(3.4), Inches(8), Inches(0.7), CARD_BG)
add_text(slide, Inches(2.7), Inches(3.45), Inches(7.6), Inches(0.6),
         "|\u03c8\u27e9 = \u2211 c\u1d62|i\u27e9      (2\u207f complex amplitudes,  \u2211|c\u1d62|\u00b2 = 1)",
         font_size=22, color=DARK, font="Cambria Math", align=PP_ALIGN.CENTER)

add_card(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.3),
         "Simulator Implementation", [
             "State stored as numpy complex128 array of 2\u207f entries",
             "Big-endian ordering: q\u2080 = most significant bit",
             "Max 5 qubits = 32 amplitudes",
             "Auto-renormalization after every gate",
         ], accent_color=BLUE, font_size=15)

add_card(slide, Inches(7), Inches(4.5), Inches(5.5), Inches(2.3),
         "Example: 2 Qubits", [
             "|00\u27e9 \u2192 index 0,   |01\u27e9 \u2192 index 1",
             "|10\u27e9 \u2192 index 2,   |11\u27e9 \u2192 index 3",
             "Bell state: [1/\u221a2, 0, 0, 1/\u221a2]",
             "Probabilities: P(00)=0.5, P(11)=0.5",
         ], accent_color=GREEN, font_size=15)

# ════════════════════════════════════════════════════════════
# SLIDE 6: SINGLE-QUBIT GATES
# ════════════════════════════════════════════════════════════
slide = content_slide("Single-Qubit Gates (Fixed)")

gates_data = [
    ("I",  "[[1,0],[0,1]]",       "Identity \u2014 no-op"),
    ("X",  "[[0,1],[1,0]]",       "Pauli-X (NOT) \u2014 bit flip: |0\u27e9\u2194|1\u27e9"),
    ("Y",  "[[0,-i],[i,0]]",      "Pauli-Y \u2014 bit+phase flip"),
    ("Z",  "[[1,0],[0,-1]]",      "Pauli-Z \u2014 phase flip: |1\u27e9\u2192-|1\u27e9"),
    ("H",  "1/\u221a2 [[1,1],[1,-1]]", "Hadamard \u2014 superposition gate"),
    ("S",  "[[1,0],[0,i]]",       "S = \u221aZ = P(\u03c0/2)"),
    ("S\u2020", "[[1,0],[0,-i]]",     "S-dagger = P(-\u03c0/2)"),
    ("T",  "[[1,0],[0,e^(i\u03c0/4)]]", "T = \u221aS = P(\u03c0/4)"),
    ("T\u2020", "[[1,0],[0,e^(-i\u03c0/4)]]", "T-dagger = P(-\u03c0/4)"),
]

y = Inches(1.25)
for name, matrix, desc in gates_data:
    # Gate name badge
    badge = add_shape(slide, Inches(0.8), y, Inches(0.7), Inches(0.42), BLUE)
    badge.text_frame.paragraphs[0].text = name
    badge.text_frame.paragraphs[0].font.size = Pt(14)
    badge.text_frame.paragraphs[0].font.color.rgb = WHITE
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.name = "Consolas"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, Inches(1.7), y, Inches(3.5), Inches(0.42),
             matrix, font_size=13, color=DARK, font="Consolas")
    add_text(slide, Inches(5.3), y, Inches(7.5), Inches(0.42),
             desc, font_size=14, color=DARK)
    y += Inches(0.52)

# Key identities box
add_card(slide, Inches(0.8), y + Inches(0.15), Inches(11.7), Inches(1.2),
         "Key Identities (verified by unit tests)", [
             "XX = YY = ZZ = HH = I     |     HXH = Z,  HZH = X     |     SS\u2020 = TT\u2020 = I     |     S\u00b2 = Z,  T\u00b2 = S",
         ], accent_color=ACCENT, font_size=15)

# ════════════════════════════════════════════════════════════
# SLIDE 7: PARAMETERIZED GATES
# ════════════════════════════════════════════════════════════
slide = content_slide("Parameterized Rotation Gates")

pgates = [
    ("P(\u03c6)",  "[[1, 0], [0, e^(i\u03c6)]]",
     "General phase gate. Applies e^(i\u03c6) to |1\u27e9.\nSpecial cases: P(\u03c0/2)=S, P(\u03c0/4)=T, P(\u03c0)=Z"),
    ("Rx(\u03b8)", "[[cos\u03b8/2, -i\u00b7sin\u03b8/2],\n [-i\u00b7sin\u03b8/2, cos\u03b8/2]]",
     "Rotation around X-axis by \u03b8 radians.\nTraces a great circle through |0\u27e9 and |1\u27e9 on the Bloch sphere."),
    ("Ry(\u03b8)", "[[cos\u03b8/2, -sin\u03b8/2],\n [sin\u03b8/2, cos\u03b8/2]]",
     "Rotation around Y-axis by \u03b8 radians.\nOnly real entries \u2014 rotates between |0\u27e9 and |1\u27e9 without phases."),
    ("Rz(\u03b8)", "[[e^(-i\u03b8/2), 0],\n [0, e^(i\u03b8/2)]]",
     "Rotation around Z-axis by \u03b8 radians.\nBloch vector rotates in the XY equatorial plane."),
]

for i, (name, matrix, desc) in enumerate(pgates):
    left = Inches(0.8) if i < 2 else Inches(6.8)
    top = Inches(1.4) + Inches(2.7) * (i % 2)
    card_w = Inches(5.5)
    card_h = Inches(2.4)

    add_shape(slide, left, top, card_w, card_h, CARD_BG)
    add_shape(slide, left, top, Inches(0.06), card_h, BLUE)

    # Gate name
    badge = add_shape(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.9), Inches(0.4), BLUE)
    badge.text_frame.paragraphs[0].text = name
    badge.text_frame.paragraphs[0].font.size = Pt(16)
    badge.text_frame.paragraphs[0].font.color.rgb = WHITE
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.name = "Consolas"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, left + Inches(1.3), top + Inches(0.15), Inches(3.8), Inches(0.8),
             matrix, font_size=13, color=DARK, font="Consolas")
    add_text(slide, left + Inches(0.2), top + Inches(1.0), Inches(5), Inches(1.2),
             desc, font_size=14, color=DARK)

# ════════════════════════════════════════════════════════════
# SLIDE 8: MULTI-QUBIT GATES
# ════════════════════════════════════════════════════════════
slide = content_slide("Multi-Qubit Gates & Entanglement")

mgates = [
    ("CNOT (CX)", "2 qubits",
     "Controlled-NOT: flips target if control = |1\u27e9.\n"
     "The primary entangling gate.\n"
     "H(q0) + CNOT(q0,q1) \u2192 Bell state (|00\u27e9+|11\u27e9)/\u221a2\n"
     "Matrix: [[1,0,0,0],[0,1,0,0],[0,0,0,1],[0,0,1,0]]"),
    ("CZ", "2 qubits",
     "Controlled-Z: applies Z to target if control = |1\u27e9.\n"
     "Symmetric: CZ(a,b) = CZ(b,a).\n"
     "Matrix: diag(1, 1, 1, -1)"),
    ("SWAP", "2 qubits",
     "Exchanges two qubit states: |01\u27e9 \u2194 |10\u27e9.\n"
     "Decomposable into 3 CNOT gates.\n"
     "Matrix swaps rows/cols for |01\u27e9 and |10\u27e9"),
    ("Toffoli (CCX)", "3 qubits",
     "Flips target only when BOTH controls = |1\u27e9.\n"
     "Universal for classical computation.\n"
     "|110\u27e9 \u2194 |111\u27e9 (8\u00d78 matrix)"),
]

for i, (name, qubits, desc) in enumerate(mgates):
    left = Inches(0.6) + Inches(3.1) * i
    add_shape(slide, left, Inches(1.3), Inches(2.9), Inches(5.5), CARD_BG)
    add_shape(slide, left, Inches(1.3), Inches(2.9), Inches(0.06), BLUE)

    add_text(slide, left + Inches(0.15), Inches(1.5), Inches(2.6), Inches(0.5),
             name, font_size=18, color=DARK_BLUE, bold=True)
    add_text(slide, left + Inches(0.15), Inches(1.95), Inches(2.6), Inches(0.3),
             qubits, font_size=13, color=ACCENT, italic=True)
    add_text(slide, left + Inches(0.15), Inches(2.4), Inches(2.6), Inches(4),
             desc, font_size=13, color=DARK)

# ════════════════════════════════════════════════════════════
# SLIDE 9: MEASUREMENT
# ════════════════════════════════════════════════════════════
section_divider("Measurement & Collapse",
                "The Born Rule & Projective Measurement")

slide = content_slide("Quantum Measurement")

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5),
         "The Born Rule", [
             "Probability of outcome i:  P(i) = |\u27e8i|\u03c8\u27e9|\u00b2 = |c\u1d62|\u00b2",
             "Post-measurement state:  |\u03c8'\u27e9 = P\u1d62|\u03c8\u27e9 / \u221aP(i)",
             "Measurement is irreversible \u2014 information is lost",
             "The outcome is fundamentally random",
         ], accent_color=BLUE, font_size=15)

add_card(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(2.5),
         "Simulator Implementation", [
             "Calculate P(0) by summing |c\u1d62|\u00b2 for basis states with bit j = 0",
             "Draw random r \u2208 [0,1): outcome = 0 if r < P(0), else 1",
             "Project: zero out amplitudes inconsistent with outcome",
             "Renormalize: divide by \u221aP(outcome)",
         ], accent_color=ACCENT, font_size=15)

add_card(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5),
         "Measurement Bases (used in BB84)", [
             "Z-basis (computational):  |0\u27e9 / |1\u27e9  \u2014  standard measurement, no transformation needed",
             "X-basis (Hadamard):  |+\u27e9 / |\u2212\u27e9  \u2014  apply H before measuring in Z to rotate |+\u27e9\u2192|0\u27e9, |\u2212\u27e9\u2192|1\u27e9",
             "Y-basis:  |+i\u27e9 / |\u2212i\u27e9  \u2014  apply S\u2020 then H before measuring in Z",
             "After basis transformation, a standard Z-measurement gives the result in the chosen basis",
         ], accent_color=GREEN, font_size=15)

# ════════════════════════════════════════════════════════════
# SLIDE 10: BLOCH SPHERE
# ════════════════════════════════════════════════════════════
slide = content_slide("Bloch Sphere Visualization")

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3.0),
         "Bloch Vector from Density Matrix", [
             "Reduced density matrix: \u03c1 = Tr_B(|\u03c8\u27e9\u27e8\u03c8|)",
             "x = 2\u00b7Re(\u03c1\u2080\u2081)  =  Tr(\u03c1\u00b7\u03c3\u2093)",
             "y = -2\u00b7Im(\u03c1\u2080\u2081) =  Tr(\u03c1\u00b7\u03c3\u1d67)",
             "z = \u03c1\u2080\u2080 - \u03c1\u2081\u2081    =  Tr(\u03c1\u00b7\u03c3\u2094)",
             "Purity r = \u221a(x\u00b2+y\u00b2+z\u00b2): 1.0 = pure, <1.0 = mixed",
         ], accent_color=BLUE, font_size=15)

add_card(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(3.0),
         "Key Positions on the Sphere", [
             "|0\u27e9  \u2192  North pole  (0, 0, +1)",
             "|1\u27e9  \u2192  South pole  (0, 0, \u22121)",
             "|+\u27e9  \u2192  +X equator  (+1, 0, 0)",
             "|\u2212\u27e9  \u2192  \u2212X equator  (\u22121, 0, 0)",
             "|+i\u27e9 \u2192  +Y equator  (0, +1, 0)",
             "|\u2212i\u27e9 \u2192  \u2212Y equator  (0, \u22121, 0)",
         ], accent_color=GREEN, font_size=15)

add_card(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.2),
         "What the Bloch Sphere Tells You", [
             "Pure qubit: vector on the surface (r = 1). Gate rotations trace paths on the sphere.",
             "Entangled qubit: vector inside the sphere (r < 1). Information is stored in correlations, not locally.",
             "Hadamard: north pole \u2192 equator (+X). Pauli-X: north \u2192 south (180\u00b0 around X-axis).",
             "Rotation gates Rx/Ry/Rz: continuous rotation around the named axis by angle \u03b8.",
         ], accent_color=ACCENT, font_size=15)

# ════════════════════════════════════════════════════════════
# SLIDE 11: ENTANGLEMENT & BELL STATES
# ════════════════════════════════════════════════════════════
slide = content_slide("Entanglement & Bell States")

add_text(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.8),
         "Two qubits are entangled when their joint state cannot be written as a tensor product |a\u27e9\u2297|b\u27e9.\n"
         "The four Bell states form a maximally entangled basis for two qubits:",
         font_size=17, color=DARK)

bell_states = [
    ("\u03a6\u207a", "(|00\u27e9 + |11\u27e9) / \u221a2", "H(q0), CNOT(q0,q1)"),
    ("\u03a6\u207b", "(|00\u27e9 \u2212 |11\u27e9) / \u221a2", "H(q0), Z(q0), CNOT(q0,q1)"),
    ("\u03a8\u207a", "(|01\u27e9 + |10\u27e9) / \u221a2", "H(q0), CNOT(q0,q1), X(q1)"),
    ("\u03a8\u207b", "(|01\u27e9 \u2212 |10\u27e9) / \u221a2", "H(q0), Z(q0), CNOT(q0,q1), X(q1)"),
]

for i, (name, state, circuit) in enumerate(bell_states):
    left = Inches(0.6) + Inches(3.1) * i
    add_shape(slide, left, Inches(2.3), Inches(2.9), Inches(2.5), CARD_BG)
    add_shape(slide, left, Inches(2.3), Inches(2.9), Inches(0.06), BLUE)
    add_text(slide, left + Inches(0.2), Inches(2.5), Inches(2.5), Inches(0.45),
             f"|{name}\u27e9", font_size=24, color=DARK_BLUE, bold=True,
             align=PP_ALIGN.CENTER, font="Cambria Math")
    add_text(slide, left + Inches(0.2), Inches(3.0), Inches(2.5), Inches(0.5),
             state, font_size=16, color=DARK, align=PP_ALIGN.CENTER, font="Cambria Math")
    add_text(slide, left + Inches(0.2), Inches(3.6), Inches(2.5), Inches(0.8),
             f"Circuit:\n{circuit}", font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

add_card(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.8),
         "No-Cloning Theorem", [
             "It is physically impossible to copy an arbitrary unknown quantum state.",
             "This is why teleportation destroys the original, BB84 is secure, and quantum info can't be broadcast.",
             "Proof: if U|0\u27e9|\u03c8\u27e9 = |\u03c8\u27e9|\u03c8\u27e9 for all |\u03c8\u27e9, then U is not unitary \u2014 contradiction.",
         ], accent_color=RED, font_size=15)

# ════════════════════════════════════════════════════════════
# SLIDE 12: DEUTSCH-JOZSA
# ════════════════════════════════════════════════════════════
section_divider("Deutsch-Jozsa Algorithm",
                "Exponential speedup for oracle classification")

slide = content_slide("Deutsch-Jozsa Algorithm")

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.3),
         "The Problem", [
             "Given f: {0,1}\u207f \u2192 {0,1}, promised to be constant or balanced",
             "Classical: needs 2\u207f\u207b\u00b9 + 1 queries (worst case)",
             "Quantum: exactly 1 query!",
             "Exponential speedup via quantum interference",
         ], accent_color=BLUE, font_size=15)

add_card(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(2.3),
         "Circuit Steps", [
             "1. Prepare |0\u27e9\u207f|1\u27e9 (n input qubits + 1 ancilla)",
             "2. Apply H to all qubits (uniform superposition)",
             "3. Apply oracle U_f (phase kickback: |x\u27e9 \u2192 (-1)^f(x)|x\u27e9)",
             "4. Apply H to input qubits",
             "5. Measure: all |0\u27e9 = constant, any |1\u27e9 = balanced",
         ], accent_color=GREEN, font_size=15)

add_card(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(3),
         "Physics: Why It Works \u2014 Quantum Interference", [
             "The ancilla in |\u2212\u27e9 enables phase kickback: the oracle encodes f(x) as a \u00b11 phase on each |x\u27e9.",
             "For constant f: all amplitudes get the same phase \u2192 final Hadamard recombines them to |0\u27e9\u207f (constructive).",
             "For balanced f: half get +1, half get \u22121 \u2192 final Hadamard causes destructive interference at |0\u27e9\u207f.",
             "Result: measurement distinguishes constant vs balanced with 100% certainty in a single query.",
             "Simulator oracle types: constant_0 (identity), constant_1 (X on ancilla), balanced (f(x) = popcount(x & pattern) mod 2).",
         ], accent_color=ACCENT, font_size=15)

# ════════════════════════════════════════════════════════════
# SLIDE 13: GROVER'S
# ════════════════════════════════════════════════════════════
section_divider("Grover\u2019s Search Algorithm",
                "Quadratic speedup for unstructured search")

slide = content_slide("Grover\u2019s Search Algorithm")

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2),
         "The Problem", [
             "Search unsorted database of N = 2\u207f items",
             "Classical: O(N) queries",
             "Quantum: O(\u221aN) queries \u2014 quadratic speedup",
         ], accent_color=BLUE, font_size=15)

add_card(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(2),
         "Optimal Iterations", [
             "k = \u230a\u03c0/4 \u00b7 \u221a(N/M)\u230b",
             "N = search space size (2\u207f), M = number of targets",
             "Too many iterations \u2192 probability decreases (overshooting)",
         ], accent_color=ORANGE, font_size=15)

add_card(slide, Inches(0.8), Inches(3.6), Inches(11.7), Inches(3.5),
         "Amplitude Amplification \u2014 The Core Physics", [
             "1. Initialize: H\u2297\u207f|0\u27e9\u207f = |s\u27e9 (uniform superposition, all amplitudes = 1/\u221aN)",
             "2. Oracle: flip phase of target states \u2014 |t\u27e9 \u2192 \u2212|t\u27e9  (diagonal matrix with \u22121 at target indices)",
             "3. Diffusion: reflect about the mean \u2014 operator 2|s\u27e9\u27e8s| \u2212 I (inversion about average amplitude)",
             "4. Geometric picture: state vector lives in 2D subspace (|target\u27e9, |non-target\u27e9). Each iteration rotates by 2\u03b8 where sin\u03b8 = \u221a(M/N).",
             "5. After ~\u03c0/(4\u03b8) iterations, state aligns with target subspace \u2192 measure with probability \u2248 1.",
             "In the simulator: watch probability histogram \u2014 target bar grows with each iteration while others shrink.",
         ], accent_color=ACCENT, font_size=15)

# ════════════════════════════════════════════════════════════
# SLIDE 14: TELEPORTATION
# ════════════════════════════════════════════════════════════
section_divider("Quantum Teleportation",
                "Transferring quantum states using entanglement + classical communication")

slide = content_slide("Quantum Teleportation Protocol")

steps = [
    ("1. Bell Pair", "H(q1), CNOT(q1,q2)\nq1 and q2 now share\n|\u03a6+\u27e9 = (|00\u27e9+|11\u27e9)/\u221a2"),
    ("2. Alice Acts", "CNOT(q0,q1), H(q0)\nEntangles message qubit\nwith her Bell half"),
    ("3. Alice Measures", "Measure q0 and q1\nObtains 2 classical bits\n(m0, m1)"),
    ("4. Bob Corrects", "If m1=1: apply X\nIf m0=1: apply Z\nq2 is now in state |\u03c8\u27e9"),
]

for i, (title_t, desc) in enumerate(steps):
    left = Inches(0.5) + Inches(3.15) * i
    add_shape(slide, left, Inches(1.3), Inches(2.9), Inches(2.3), CARD_BG)
    add_shape(slide, left, Inches(1.3), Inches(2.9), Inches(0.06), BLUE)

    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.05), Inches(1.5), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = str(i + 1)
    circle.text_frame.paragraphs[0].font.size = Pt(20)
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, left + Inches(0.1), Inches(2.1), Inches(2.7), Inches(0.4),
             title_t, font_size=15, color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), Inches(2.5), Inches(2.7), Inches(1),
             desc, font_size=13, color=DARK, align=PP_ALIGN.CENTER)

add_card(slide, Inches(0.8), Inches(3.9), Inches(11.7), Inches(3),
         "Physics: Why Teleportation Works", [
             "After Alice's CNOT+H, the 3-qubit state decomposes into 4 terms conditioned on Alice's measurement outcomes.",
             "Each outcome leaves Bob's qubit in a known Pauli rotation of the original |\u03c8\u27e9 = \u03b1|0\u27e9 + \u03b2|1\u27e9.",
             "Classical bits (m0, m1) tell Bob which correction to apply: I, X, Z, or XZ.",
             "No faster-than-light communication: Bob's qubit is random until he receives Alice's classical bits.",
             "The original state is destroyed (no-cloning): Alice's qubit collapses upon measurement.",
             "Simulator: set \u03b8/\u03c6 sliders for q0, run protocol, verify Bob's Bloch sphere (q2) matches the original.",
         ], accent_color=ACCENT, font_size=15)

# ════════════════════════════════════════════════════════════
# SLIDE 15: BB84
# ════════════════════════════════════════════════════════════
section_divider("BB84 Quantum Key Distribution",
                "Provably secure cryptography from quantum mechanics")

slide = content_slide("BB84 Protocol")

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3),
         "Protocol Steps", [
             "1. Alice generates random bits, picks random basis (Z or X) each",
             "2. Z-basis: 0\u2192|0\u27e9, 1\u2192|1\u27e9.  X-basis: 0\u2192|+\u27e9, 1\u2192|\u2212\u27e9",
             "3. Bob picks random measurement basis per qubit",
             "4. Sifting: keep bits where bases matched (~50%)",
             "5. Estimate QBER on a sacrificed subset",
             "6. QBER < 11% \u2192 secure.  QBER \u2265 11% \u2192 Eve detected!",
         ], accent_color=BLUE, font_size=15)

add_card(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(3),
         "Why Eve Gets Caught", [
             "Eve must measure to read \u2014 but doesn't know Alice's basis",
             "Wrong basis guess (50%): measurement collapses to wrong basis",
             "Re-sent qubit: Bob gets error 50% of the time (when matching Alice)",
             "Net effect: Eve introduces ~25% QBER",
             "Security threshold: 11% (provably secure bound for BB84)",
         ], accent_color=RED, font_size=15)

add_card(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.3),
         "Security Foundation \u2014 Quantum Physics Guarantees", [
             "No-Cloning Theorem: Eve cannot copy qubits to measure in both bases simultaneously.",
             "Measurement Disturbance: any measurement on a qubit in a non-matching basis irreversibly disturbs it.",
             "Information-theoretic security: security does not depend on computational assumptions \u2014 it comes from physics.",
             "Simulator: toggle Eve on/off and observe QBER jump from ~0% to ~25%.",
         ], accent_color=ACCENT, font_size=15)

# ════════════════════════════════════════════════════════════
# SLIDE 16: QRNG
# ════════════════════════════════════════════════════════════
slide = content_slide("Quantum Random Number Generator (QRNG)")

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5),
         "The Concept", [
             "Classical PRNGs are deterministic (predictable with seed)",
             "Quantum mechanics provides TRUE randomness",
             "Measuring a superposition is fundamentally unpredictable",
             "No hidden variables (Bell's theorem)",
         ], accent_color=BLUE, font_size=16)

add_card(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(2.5),
         "The Circuit (per bit)", [
             "1. Prepare fresh qubit: |0\u27e9",
             "2. Apply Hadamard: |0\u27e9 \u2192 (|0\u27e9+|1\u27e9)/\u221a2",
             "3. Measure: exactly 50/50 probability",
             "4. Repeat N times for N truly random bits",
         ], accent_color=GREEN, font_size=16)

add_card(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5),
         "Simulator Output", [
             "Bitstring: the generated random binary sequence (up to 256 bits)",
             "Integer / Hex value: decimal and hexadecimal representation (\u226464 bits)",
             "Frequency analysis: count of 0s vs 1s \u2014 should approach 50/50 for large N",
             "State history: step-by-step statevector at each H \u2192 Measure cycle",
         ], accent_color=ACCENT, font_size=16)

# ════════════════════════════════════════════════════════════
# SLIDE 17: DEMO WALKTHROUGH
# ════════════════════════════════════════════════════════════
slide = content_slide("Live Demo Walkthrough")

demos = [
    ("Demo 1: Superposition", [
        "1 qubit system",
        "Apply H to q0",
        "Observe Bloch sphere: north pole \u2192 equator",
        "Probability chart: 50/50 split",
    ]),
    ("Demo 2: Bell State", [
        "2 qubit system",
        "Apply H(q0), then CNOT(q0,q1)",
        "Both Bloch spheres shrink (entangled)",
        "Probabilities: only |00\u27e9 and |11\u27e9",
    ]),
    ("Demo 3: Grover\u2019s Search", [
        "3 qubits, target = 5 (|101\u27e9)",
        "Watch amplitude amplification",
        "Target probability rises each iteration",
        "Final measurement finds target",
    ]),
    ("Demo 4: Teleportation", [
        "Set arbitrary state on q0",
        "Run protocol, step through",
        "Observe Bob\u2019s q2 matches original",
        "Original q0 is destroyed (no-cloning)",
    ]),
]

for i, (demo_title, steps) in enumerate(demos):
    left = Inches(0.5) + Inches(3.15) * i
    add_card(slide, left, Inches(1.3), Inches(2.9), Inches(3),
             demo_title, steps,
             accent_color=[BLUE, GREEN, ACCENT, ORANGE][i], font_size=14)

add_card(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.2),
         "Tips for Exploration", [
             "Use Step Forward/Backward to freeze time and inspect states between gates.",
             "Compare Deutsch-Jozsa with constant_0 vs balanced \u2014 see the interference patterns differ.",
             "Toggle Eve in BB84: watch QBER jump from ~0% to ~25%.",
             "Try QRNG with 200 bits and verify frequency approaches 50/50.",
         ], accent_color=DARK_BLUE, font_size=15)

# ════════════════════════════════════════════════════════════
# SLIDE 18: THANK YOU
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_text(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
         "Thank You",
         font_size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.04), BLUE)
add_text(slide, Inches(1), Inches(3.9), Inches(11.3), Inches(1),
         "Quantum Protocol & Algorithm Simulator",
         font_size=24, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.7), Inches(11.3), Inches(1),
         "Questions?",
         font_size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# ── Save ──
output = "Quantum_Simulator_Presentation.pptx"
prs.save(output)
print(f"Created: {output}")
